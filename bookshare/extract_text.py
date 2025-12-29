from pptx import Presentation
import json
import sys
from pathlib import Path

def extract_text(pptx_path: str):
    prs = Presentation(pptx_path)
    out = {"source": str(pptx_path), "slides": []}

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)

        out["slides"].append({
            "slide_index": i,
            "texts": texts
        })

    return out

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_text.py path/to/file.pptx")
        sys.exit(1)

    # ファイルパスにスペースが含まれる場合、複数の引数に分割される可能性があるため結合
    pptx_path = " ".join(sys.argv[1:])
    data = extract_text(pptx_path)

    out_path = Path(pptx_path).with_suffix(".texts.json")
    out_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Saved: {out_path}")
