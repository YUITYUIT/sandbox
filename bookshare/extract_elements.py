from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import sys
from pathlib import Path

EMU_PER_INCH = 914400
CM_PER_INCH = 2.54

def emu_to_cm(emu: int) -> float:
    return round((emu / EMU_PER_INCH) * CM_PER_INCH, 3)

def safe_int(v):
    try:
        return int(v)
    except Exception:
        return None

def extract_slide_elements(pptx_path: str):
    prs = Presentation(pptx_path)
    deck = {
        "source": str(pptx_path),
        "slide_size_cm": {
            "width": emu_to_cm(prs.slide_width),
            "height": emu_to_cm(prs.slide_height),
        },
        "slides": []
    }

    for si, slide in enumerate(prs.slides, start=1):
        slide_obj = {"slide_index": si, "elements": []}

        for idx, shape in enumerate(slide.shapes, start=1):
            stype = shape.shape_type  # enum MSO_SHAPE_TYPE
            el = {
                "id_in_slide": idx,
                "name": getattr(shape, "name", None),
                "shape_type": str(stype),
                "kind": None,
                "bbox_cm": None,
                "text": None,
                "text_lines": None,
                "image": None,
                "table": None,
            }

            # bbox (position & size)
            # Some shapes might miss these; guard
            if hasattr(shape, "left") and hasattr(shape, "top") and hasattr(shape, "width") and hasattr(shape, "height"):
                el["bbox_cm"] = {
                    "x": emu_to_cm(shape.left),
                    "y": emu_to_cm(shape.top),
                    "w": emu_to_cm(shape.width),
                    "h": emu_to_cm(shape.height),
                }

            # classify & extract
            if stype == MSO_SHAPE_TYPE.PICTURE:
                el["kind"] = "image"
                # image info is limited without saving; but we can keep placeholder
                el["image"] = {"exists": True}

            elif stype == MSO_SHAPE_TYPE.TABLE:
                el["kind"] = "table"
                tbl = shape.table
                rows, cols = len(tbl.rows), len(tbl.columns)
                cells = []
                for r in range(rows):
                    row = []
                    for c in range(cols):
                        txt = tbl.cell(r, c).text.strip()
                        row.append(txt)
                    cells.append(row)
                el["table"] = {
                    "rows": rows,
                    "cols": cols,
                    "cells": cells
                }

            elif getattr(shape, "has_chart", False):
                el["kind"] = "chart"
                el["chart"] = {"exists": True}

            elif getattr(shape, "has_text_frame", False):
                el["kind"] = "text"
                tf = shape.text_frame
                txt = tf.text.strip()
                el["text"] = txt if txt else None
                if txt:
                    el["text_lines"] = [line.strip() for line in txt.splitlines() if line.strip()]

            else:
                # shapes, groups, smartart, etc.
                el["kind"] = "other"

            # Drop empty noise elements (optional): keep all for now
            slide_obj["elements"].append(el)

        deck["slides"].append(slide_obj)

    return deck

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_elements.py path/to/file.pptx")
        sys.exit(1)

    # ファイルパスにスペースが含まれる場合、複数の引数に分割される可能性があるため結合
    pptx_path = " ".join(sys.argv[1:])
    data = extract_slide_elements(pptx_path)

    out_path = Path(pptx_path).with_suffix(".elements.json")
    out_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Saved: {out_path}")
